from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
import os
import win32com.client
from pathlib import Path
from PIL import Image
import io
import tempfile
from ..images.image_processor import ImageProcessor
from .ppt_extractor import PPTExtractor
import logging

class PPTProcessor:
    """PPT处理器 - 专注于PPT编辑和格式调整功能"""
    
    def __init__(self, db_manager=None):
        self.current_ppt = None
        self.current_ppt_path = None
        self.db_manager = db_manager
        self.image_processor = ImageProcessor(db_manager) if db_manager else None
        self.ppt_extractor = PPTExtractor(db_manager) if db_manager else None
        
    def get_image_processor(self):
        """获取图片处理器实例"""
        if self.image_processor is None:
            raise RuntimeError("ImageProcessor未初始化，请确保提供了有效的db_manager")
        return self.image_processor
    
    def open_presentation(self, filepath: str):
        """打开PPT文件"""
        try:
            self.current_ppt = Presentation(filepath)
            self.current_ppt_path = filepath
        except Exception as e:
            raise ValueError(f"打开PPT文件失败: {str(e)}")
    
    def adjust_text_boxes(self):
        """调整所有文本框：设置为自动调整大小且不自动换行"""
        if not self.current_ppt_path:
            raise ValueError("未打开PPT文件")
        
        try:
            # 获取PowerPoint应用程序实例
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            
            # 打开演示文稿
            abs_path = os.path.abspath(self.current_ppt_path)
            presentation = powerpoint.Presentations.Open(abs_path)
            
            try:
                # 遍历所有幻灯片
                for slide in presentation.Slides:
                    # 遍历幻灯片中的所有形状
                    for shape in slide.Shapes:
                        if shape.HasTextFrame:
                            # 取消勾选"形状中的文字自动换行"
                            shape.TextFrame.WordWrap = False
                            
                            # 先设置为不自动调整
                            shape.TextFrame.AutoSize = 0  # ppAutoSizeNone
                            
                            # 最后设置为根据文字调整形状大小
                            shape.TextFrame.AutoSize = 1  # ppAutoSizeShapeToFitText
                
                # 保存更改
                presentation.Save()
                
            finally:
                # 关闭演示文稿
                presentation.Close()
                
        except Exception as e:
            raise Exception(f"调整文本框时出错: {str(e)}")
        
        finally:
            try:
                # 确保PowerPoint应用程序被正确关闭
                powerpoint.Quit()
            except:
                pass
    
    def extract_text(self) -> str:
        """提取PPT中的所有文字内容"""
        if not self.current_ppt:
            raise ValueError("未打开PPT文件")
            
        text_content = []
        
        # 遍历所有幻灯片
        for slide_index, slide in enumerate(self.current_ppt.slides, 1):
            text_content.append(f"\n--- 第{slide_index}页 ---\n")
            
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # 获取形状的文本内容
                    text = shape.text.strip()
                    if text:
                        text_content.append(text)
        
        # 合并所有文本
        return "\n".join(text_content)
    
    def clean_unused_layouts(self) -> int:
        """彻底清理PPT母版"""
        if not self.current_ppt_path:
            raise ValueError("未打开PPT文件")
        
        try:
            cleaned_count = 0
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            abs_path = os.path.abspath(self.current_ppt_path)
            presentation = powerpoint.Presentations.Open(abs_path)
            
            try:
                # 收集所有使用中的设计和布局名称
                used_designs = set()
                used_layouts = set()
                
                for slide in presentation.Slides:
                    if hasattr(slide, 'Design') and slide.Design:
                        used_designs.add(slide.Design.Name)
                    if hasattr(slide, 'CustomLayout') and slide.CustomLayout:
                        used_layouts.add(slide.CustomLayout.Name)
                
                # 删除未使用的设计（母版）
                for design in list(presentation.Designs):
                    design_name = design.Name
                    if design_name not in used_designs:
                        try:
                            design.Delete()
                            cleaned_count += 1
                            print(f"删除未使用的设计: {design_name}")
                        except Exception as e:
                            print(f"删除设计 {design_name} 时出错: {str(e)}")
                
                # 处理剩余的设计（母版）
                for design in presentation.Designs:
                    slide_master = design.SlideMaster
                    
                    # 删除未使用的布局
                    for layout in list(slide_master.CustomLayouts):
                        layout_name = layout.Name
                        if layout_name not in used_layouts:
                            try:
                                layout.Delete()
                                cleaned_count += 1
                                print(f"删除未使用的布局: {layout_name}")
                            except Exception as e:
                                print(f"删除布局 {layout_name} 时出错: {str(e)}")
                    
                    # 删除母版中的所有占位符
                    for shape in list(slide_master.Shapes):
                        try:
                            if shape.Type == 14:  # 14 表示占位符
                                shape.Delete()
                                cleaned_count += 1
                                print("删除母版中的占位符")
                        except Exception as e:
                            print(f"删除母版占位符时出错: {str(e)}")
                    
                    # 删除所有布局中的占位符
                    for layout in slide_master.CustomLayouts:
                        for shape in list(layout.Shapes):
                            try:
                                if shape.Type == 14:  # 14 表示占位符
                                    shape.Delete()
                                    cleaned_count += 1
                                    print(f"删除布局 {layout.Name} 中的占位符")
                            except Exception as e:
                                print(f"删除布局占位符时出错: {str(e)}")
                
                # 保存更改
                presentation.Save()
                print(f"清理完成，共删除 {cleaned_count} 个元素")
                
            finally:
                # 关闭演示文稿
                presentation.Close()
                
        except Exception as e:
            raise Exception(f"清理模板时出错: {str(e)}")
        
        finally:
            try:
                # 确保PowerPoint应用程序被正确关闭
                powerpoint.Quit()
            except:
                pass
        
        return cleaned_count    
    def save(self, filepath: str = None):
        """保存PPT文件"""
        if not self.current_ppt:
            raise ValueError("未打开PPT文件")
            
        # 如果没有指定保存路径，在原文件名后添加后缀
        if not filepath:
            base, ext = os.path.splitext(self.current_ppt_path)
            filepath = f"{base}_modified{ext}"
        
        self.current_ppt.save(filepath)
    
    def extract_all_images(self, output_folder: str) -> list:
        """从PPT中提取所有图片（保持原始格式）"""
        if not self.current_ppt:
            raise ValueError("未打开PPT文件")
        
        try:
            output_folder = Path(output_folder)
            output_folder.mkdir(parents=True, exist_ok=True)
            extracted_images = []
            
            # 遍历所有幻灯片
            for slide_idx, slide in enumerate(self.current_ppt.slides, 1):
                # 提取形状中的图片
                for shape_idx, shape in enumerate(slide.shapes, 1):
                    try:
                        if hasattr(shape, 'image'):
                            # 获取图片数据
                            img_data = shape.image.blob
                            content_type = shape.image.content_type
                            
                            # 确定文件扩展名
                            ext = self._get_image_extension(content_type, img_data)
                            
                            # 生成图片文件名
                            img_name = f"slide{slide_idx}_shape{shape_idx}{ext}"
                            img_path = str(output_folder / img_name)
                            
                            # 保存原始图片数据
                            with open(img_path, 'wb') as f:
                                f.write(img_data)
                            
                            # 记录图片信息
                            image_info = {
                                'path': img_path,
                                'slide': slide_idx,
                                'shape': shape_idx,
                                'format': ext[1:].upper()
                            }
                            
                            extracted_images.append(image_info)
                            
                    except Exception as e:
                        print(f"处理图片时出错: {str(e)}")
                        continue
                
                # 提取背景图片
                try:
                    if hasattr(slide, 'background') and hasattr(slide.background, 'image'):
                        img_data = slide.background.image.blob
                        content_type = slide.background.image.content_type
                        ext = self._get_image_extension(content_type, img_data)
                        
                        img_name = f"slide{slide_idx}_background{ext}"
                        img_path = str(output_folder / img_name)
                        
                        with open(img_path, 'wb') as f:
                            f.write(img_data)
                        
                        image_info = {
                            'path': img_path,
                            'slide': slide_idx,
                            'shape': 'background',
                            'format': ext[1:].upper()
                        }
                        
                        extracted_images.append(image_info)
                        
                except Exception as e:
                    print(f"处理背景图片时出错: {str(e)}")
            
            return extracted_images
            
        except Exception as e:
            print(f"提取图片时出错: {str(e)}")
            raise
    
    def _get_image_extension(self, content_type: str, img_data: bytes) -> str:
        """根据内容类型和文件头确定图片扩展名"""
        ext_map = {
            'image/png': '.png',
            'image/jpeg': '.jpg',
            'image/gif': '.gif',
            'image/bmp': '.bmp',
            'image/tiff': '.tiff',
            'image/x-emf': '.emf',
            'image/x-wmf': '.wmf',
            'image/webp': '.webp',
            'image/x-icon': '.ico',
            'image/svg+xml': '.svg'
        }
        
        if content_type in ext_map:
            return ext_map[content_type]
        
        # 从文件头判断
        if img_data.startswith(b'\x89PNG'):
            return '.png'
        elif img_data.startswith(b'\xFF\xD8'):
            return '.jpg'
        elif img_data.startswith(b'GIF8'):
            return '.gif'
        elif img_data.startswith(b'BM'):
            return '.bmp'
        elif img_data.startswith(b'II*\x00') or img_data.startswith(b'MM\x00*'):
            return '.tiff'
        elif img_data.startswith(b'%PDF'):
            return '.pdf'
        
        return '.bin' 
    
    def remove_ppt_source(self, path: str):
        """移除PPT源文件
        
        Args:
            path: PPT文件路径
        """
        try:
            if not self.db_manager:
                raise RuntimeError("数据库管理器未初始化")
            
            # 从数据库中移除PPT源记录
            self.db_manager.execute(
                "DELETE FROM ppt_sources WHERE path = ?",
                (path,)
            )
            
            # 移除相关的图片映射
            self.db_manager.execute(
                "DELETE FROM image_ppt_mapping WHERE pptx_path = ?",
                (path,)
            )
            
            # 提交事务
            self.db_manager.commit()
            
            logging.info(f"成功移除PPT源文件: {path}")
            
        except Exception as e:
            logging.error(f"移除PPT源文件失败: {str(e)}")
            if self.db_manager:
                self.db_manager.rollback()
            raise
