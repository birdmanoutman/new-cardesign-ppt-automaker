import io
import logging
import os
import sqlite3
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import hashlib
from datetime import datetime
from PIL import ImageDraw, ImageFont
from functools import lru_cache
from typing import List, Dict, Tuple
import numpy as np

# 延迟导入CLIP相关库
CLIP_AVAILABLE = False
try:
    import torch
    if not torch.cuda.is_available():
        print("警告: 未检测到GPU，CLIP模型将在CPU上运行，可能会比较慢")
    
    from transformers import CLIPProcessor, CLIPModel
    CLIP_AVAILABLE = True
    print("成功加载CLIP模型相关库")
except ImportError as e:
    print(f"CLIP模型相关库导入失败: {str(e)}")
    print("请确保已正确安装以下库:")
    print("- torch")
    print("- transformers")
    print("可以尝试运行: pip install torch transformers --index-url https://pypi.tuna.tsinghua.edu.cn/simple")
except Exception as e:
    print(f"初始化CLIP相关功能时出错: {str(e)}")

class ImageProcessor:
    def __init__(self):
        self.db_path = None
        self.db_conn = None
        self.cursor = None
        self.table_name = "images"
        
        # 设置数据库路径（在应用程序数据目录）
        app_data_dir = Path(self._get_app_data_dir())
        self.db_path = app_data_dir / "image_gallery.db"
        
        # 确保目录存在
        self.db_path.parent.mkdir(parents=True, exist_ok=True)
        
        # 初始化数据库
        self.init_database()
        
        # CLIP模型相关
        self.clip_model = None
        self.clip_processor = None
        self.clip_available = CLIP_AVAILABLE
        
        if self.clip_available:
            print("CLIP功能已启用")
        else:
            print("CLIP功能未启用，标签识别将不可用")
    
    def _get_app_data_dir(self) -> Path:
        """获取应用程序数据目录"""
        if os.name == 'nt':  # Windows
            app_data = os.getenv('APPDATA')
            return Path(app_data) / 'CarDesignTools'
        else:  # Linux/Mac
            home = os.path.expanduser('~')
            return Path(home) / '.cardesigntools'
    
    def init_database(self):
        """初始化数据库"""
        try:
            # 连接数据库
            self.db_conn = sqlite3.connect(str(self.db_path))
            self.cursor = self.db_conn.cursor()
            
            # 创建表
            self._create_table()
            
            print(f"数据库初始化成功: {self.db_path}")
            
        except Exception as e:
            print(f"初始化数据库时出错: {str(e)}")
            raise
    
    def _create_table(self):
        """创建数据库表"""
        # 图片主表
        table_sql = f'''
            CREATE TABLE IF NOT EXISTS {self.table_name} (
                img_hash TEXT PRIMARY KEY,  -- 使用哈希作为主键
                img_path TEXT,              -- 图片保存路径
                img_name TEXT,              -- 图片文件名
                extract_date TEXT,          -- 提取时间
                img_type TEXT,              -- 图片类型
                format TEXT,                -- 图片格式
                width INTEGER,              -- 图片宽度
                height INTEGER,             -- 图片高度
                file_size INTEGER           -- 文件大小
            )
        '''
        self.cursor.execute(table_sql)
        
        # 图片-PPT映射关系表
        mapping_sql = '''
            CREATE TABLE IF NOT EXISTS image_ppt_mapping (
                img_hash TEXT,              -- 图片哈希值
                pptx_path TEXT,             -- PPT文件路径
                slide_index INTEGER,        -- 幻灯片索引
                shape_index TEXT,           -- 形状引
                extract_date TEXT,          -- 提取时间
                PRIMARY KEY (img_hash, pptx_path, slide_index, shape_index),
                FOREIGN KEY (img_hash) REFERENCES images(img_hash)
                    ON DELETE CASCADE       -- 当图片被删除时，自动删除映射
            )
        '''
        self.cursor.execute(mapping_sql)
        
        # 添加设置表
        settings_sql = '''
            CREATE TABLE IF NOT EXISTS settings (
                key TEXT PRIMARY KEY,
                value TEXT
            )
        '''
        self.cursor.execute(settings_sql)
        
        # 添加PPT源表
        sources_sql = '''
            CREATE TABLE IF NOT EXISTS ppt_sources (
                path TEXT PRIMARY KEY,
                added_date TEXT
            )
        '''
        self.cursor.execute(sources_sql)
        
        # 添加标签表
        tags_sql = '''
            CREATE TABLE IF NOT EXISTS tags (
                id INTEGER PRIMARY KEY,
                name TEXT UNIQUE,          -- 标签名称
                category TEXT,             -- 标签类别
                created_at TEXT            -- 创建时间
            )
        '''
        self.cursor.execute(tags_sql)
        
        # 添加图片-标签关联表
        image_tags_sql = '''
            CREATE TABLE IF NOT EXISTS image_tags (
                img_hash TEXT,             -- 图片哈希
                tag_id INTEGER,            -- 标签ID
                confidence REAL,           -- 置信度
                created_at TEXT,           -- 创建时间
                PRIMARY KEY (img_hash, tag_id),
                FOREIGN KEY (img_hash) REFERENCES images(img_hash) ON DELETE CASCADE,
                FOREIGN KEY (tag_id) REFERENCES tags(id) ON DELETE CASCADE
            )
        '''
        self.cursor.execute(image_tags_sql)
        
        # 添加索引
        self.cursor.execute("CREATE INDEX IF NOT EXISTS idx_image_tags_hash ON image_tags(img_hash)")
        self.cursor.execute("CREATE INDEX IF NOT EXISTS idx_image_tags_tag ON image_tags(tag_id)")
        
        # 添加索引以提高查询性能
        index_sql = [
            f"CREATE INDEX IF NOT EXISTS idx_{self.table_name}_hash ON {self.table_name}(img_hash)",
            "CREATE INDEX IF NOT EXISTS idx_mapping_hash ON image_ppt_mapping(img_hash)",
            "CREATE INDEX IF NOT EXISTS idx_mapping_path ON image_ppt_mapping(pptx_path)",
            "CREATE INDEX IF NOT EXISTS idx_extract_date ON image_ppt_mapping(extract_date)"
        ]
        
        for sql in index_sql:
            self.cursor.execute(sql)
        
        self.db_conn.commit()
    
    def _calculate_image_hash(self, image_data: bytes) -> str:
        """计算图片哈希"""
        return hashlib.md5(image_data).hexdigest()
    
    def _is_size_similar(self, img: Image.Image, presentation: Presentation) -> bool:
        """检图片尺寸是否接近PPT幻灯片尺寸"""
        # 获取PPT幻灯片尺寸（英寸转像素，假设96 DPI）
        ppt_width = int(presentation.slide_width * 96 / 914400)  # EMU到英寸到像素
        ppt_height = int(presentation.slide_height * 96 / 914400)
        
        # 获取图片尺寸
        img_width, img_height = img.size
        
        # 计比例
        width_ratio = img_width / ppt_width
        height_ratio = img_height / ppt_height
        
        # 如果图片尺寸接近或大于PPT尺寸，认为是景图
        return width_ratio >= 0.8 and height_ratio >= 0.8
    
    def _is_duplicate(self, img_hash: str, check_file_exists: bool = True) -> bool:
        """
        检查图片是否重复
        check_file_exists: 是否检查文件是否存在
        """
        self.cursor.execute(
            f"SELECT img_path FROM {self.table_name} WHERE img_hash=? LIMIT 1",
            (img_hash,)
        )
        result = self.cursor.fetchone()
        
        if not result:
            return False
        
        if check_file_exists:
            # 检查文件是否实际存在
            img_path = result[0]
            if not os.path.exists(img_path):
                # 如果文件不存在，删除数据库记录
                self.cursor.execute(
                    f"DELETE FROM {self.table_name} WHERE img_hash=?",
                    (img_hash,)
                )
                self.db_conn.commit()
                return False
        
        return True
    
    def _save_image(self, img: Image.Image, img_path: str):
        """保存图片"""
        # 确保保存路径存在
        os.makedirs(os.path.dirname(img_path), exist_ok=True)
        
        # 转换图片格式并保存
        if img.mode in ['RGBA', 'P']:
            img = img.convert('RGB')
        img.save(img_path)
    
    def _determine_image_type(self, content_type: str, img_data: bytes) -> str:
        """
        确定图片类型
        返回: 'background' | 'icon' | 'normal'
        """
        try:
            # 如果是WMF/EMF格式，标记为icon
            if content_type in ['image/x-wmf', 'image/x-emf']:
                return 'icon'
            
            # 打开图片获取尺寸
            img = Image.open(io.BytesIO(img_data))
            width, height = img.size
            
            # 如果尺寸太小，可能是图标
            if width < 100 or height < 100:
                return 'icon'
            
            # 如果是大尺寸图片，可能是背景
            if width >= 800 and height >= 600:
                return 'background'
            
            return 'normal'
        except:
            return 'normal'  # 如果无法判断，默认为普通图片
    
    def extract_background_images(self, ppt_path: str, output_folder: str, progress_callback=None) -> list:
        """从PPT中提取所图片建立索引"""
        try:
            self.db_conn.execute("BEGIN TRANSACTION")
            ppt_path = Path(ppt_path)
            output_folder = Path(output_folder)
            output_folder.mkdir(parents=True, exist_ok=True)
            
            presentation = Presentation(ppt_path)
            total_slides = len(presentation.slides)
            insert_data = []
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                if progress_callback:
                    progress_callback(slide_idx, total_slides, 
                        f"正在处理第 {slide_idx}/{total_slides} 页")
                
                for shape_idx, shape in enumerate(slide.shapes, 1):
                    try:
                        if not hasattr(shape, 'image'):
                            continue
                        
                        # 获取图片数据和格式
                        img_data = shape.image.blob
                        content_type = shape.image.content_type
                        
                        # 只处理JPG和PNG
                        if content_type not in ['image/jpeg', 'image/png']:
                            continue
                        
                        img_hash = self._calculate_image_hash(img_data)
                        
                        # 检查是否重复
                        if self._is_duplicate(img_hash):
                            self._add_ppt_mapping(img_hash, ppt_path, slide_idx, shape_idx)
                            continue
                        
                        # 直接保存原始数据
                        ext = '.jpg' if content_type == 'image/jpeg' else '.png'
                        img_name = f"{img_hash[:8]}{ext}"
                        img_path = output_folder / img_name
                        
                        # 直接写入原始数据
                        with open(img_path, 'wb') as f:
                            f.write(img_data)
                        
                        # 记录图片信息
                        image_info = {
                            'hash': img_hash,
                            'path': str(img_path),
                            'name': img_name,
                            'type': 'normal',
                            'format': ext[1:].upper(),
                            'file_size': len(img_data)
                        }
                        insert_data.append(image_info)
                        
                        # 添加PPT映射
                        self._add_ppt_mapping(img_hash, ppt_path, slide_idx, shape_idx)
                    
                    except Exception as e:
                        print(f"处理图时出错 (slide {slide_idx}, shape {shape_idx}): {str(e)}")
                        continue
            
            if insert_data:
                self._batch_insert_images(insert_data)
            
            self.db_conn.commit()
            return insert_data
            
        except Exception as e:
            self.db_conn.rollback()
            raise

    def _process_and_save_image(self, img_data, content_type, ppt_path, output_folder, 
                              slide_idx, shape_idx) -> dict:
        """处理并保存单个图片"""
        try:
            img_hash = self._calculate_image_hash(img_data)
            
            # 检查是否已存在相同哈希的图片
            self.cursor.execute(
                f"SELECT img_path FROM {self.table_name} WHERE img_hash = ?",
                (img_hash,)
            )
            existing = self.cursor.fetchone()
            
            if existing:
                # 如果图片已存在，只添加PPT映射关系
                self._add_ppt_mapping(img_hash, ppt_path, slide_idx, shape_idx)
                return None
            
            # 处理新片
            ext = self._get_image_extension(content_type, img_data)
            img_name = f"{img_hash[:8]}{ext}"  # 使用哈希值作为文件名
            img_path = output_folder / img_name
            
            # 保存图片
            try:
                img = Image.open(io.BytesIO(img_data))
                if img.format == 'MPO' or img.mode in ['RGBA', 'P']:
                    img = img.convert('RGB')
                img.save(img_path, quality=95, optimize=True)
                width, height = img.size
            except:
                img_path.parent.mkdir(parents=True, exist_ok=True)
                with open(img_path, 'wb') as f:
                    f.write(img_data)
                width = height = 0
            
            # 插入图片信息
            self.cursor.execute(
                f"""
                INSERT INTO {self.table_name} 
                (img_hash, img_path, img_name, extract_date, img_type, format, 
                width, height, file_size)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (img_hash, str(img_path), img_name, 
                 datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                 self._determine_image_type(content_type, img_data),
                 ext[1:].upper(), width, height, len(img_data))
            )
            
            # 添加PPT映射关系
            self._add_ppt_mapping(img_hash, ppt_path, ppt_path.name, slide_idx, shape_idx)
            
            return {
                'hash': img_hash,
                'path': str(img_path),
                'name': img_name
            }
            
        except Exception as e:
            print(f"处理图片失败: {str(e)}")
            return None

    def _add_ppt_mapping(self, img_hash: str, ppt_path: Path, slide_idx: int, shape_idx: str):
        """添加图片-PPT映射关系"""
        try:
            self.cursor.execute(
                """
                INSERT OR REPLACE INTO image_ppt_mapping 
                (img_hash, pptx_path, slide_index, shape_index, extract_date)
                VALUES (?, ?, ?, ?, ?)
                """,
                (img_hash, str(ppt_path), slide_idx, str(shape_idx),
                 datetime.now().strftime('%Y-%m-%d %H:%M:%S'))
            )
        except Exception as e:
            print(f"添加PPT映射关系时出错: {str(e)}")

    def _batch_insert_images(self, image_info_list: list):
        """批量插入图片信息"""
        if not image_info_list:
            return
        
        insert_sql = f"""
            INSERT OR REPLACE INTO {self.table_name} (
                img_hash, img_path, img_name, extract_date, img_type, format,
                width, height, file_size
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
        """
        
        # 准备批量插入的数据
        values = [(
            info['hash'], info['path'], info['name'],
            datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            info.get('type', 'normal'),
            info.get('format', 'UNKNOWN'),
            info.get('width', 0),
            info.get('height', 0),
            info.get('file_size', 0)
        ) for info in image_info_list]
        
        # 执行批量插入
        self.cursor.executemany(insert_sql, values)
        self.db_conn.commit()

    def _get_image_extension(self, content_type: str, img_data: bytes) -> str:
        """根据内容类型和文件头确定图片扩展名"""
        try:
            # 尝试使用PIL打开图片确定格式
            img = Image.open(io.BytesIO(img_data))
            if img.format:
                # 特殊处理 MPO 格式
                if img.format == 'MPO':
                    return '.jpg'  # MPO 实际上是多帧的 JPEG
                return f".{img.format.lower()}"
            
        except Exception:
            # 如果PIL无法打开，尝试从MIME类型和文件头判断
            ext_map = {
                'image/png': '.png',
                'image/jpeg': '.jpg',
                'image/gif': '.gif',
                'image/bmp': '.bmp',
                'image/tiff': '.tiff',
                'image/x-emf': '.emf',
                'image/x-wmf': '.wmf',
                'image/webp': '.webp',
                'image/x-icon': '.ico',
                'image/svg+xml': '.svg',
                'image/mpo': '.jpg',  # 添加 MPO 映射
            }
            
            if content_type in ext_map:
                return ext_map[content_type]
            
            # 从文件头判断
            if img_data.startswith(b'\x89PNG'):
                return '.png'
            elif img_data.startswith(b'\xFF\xD8'):
                return '.jpg'
            elif img_data.startswith(b'GIF8'):
                return '.gif'
            elif img_data.startswith(b'BM'):
                return '.bmp'
            elif img_data.startswith(b'II*\x00') or img_data.startswith(b'MM\x00*'):
                return '.tiff'
        
        # 如果无法确定格式，默认保存为JPEG
        return '.jpg'

    def get_all_images(self, offset=0, limit=None, filters=None) -> list:
        """获取所有图片信息"""
        try:
            query = f"""
                WITH image_counts AS (
                    SELECT img_hash, COUNT(DISTINCT pptx_path) as ref_count
                    FROM image_ppt_mapping
                    GROUP BY img_hash
                )
                SELECT DISTINCT 
                    i.img_hash, i.img_path, i.img_name, i.extract_date, 
                    i.img_type, i.width, i.height,
                    m.pptx_path,
                    COALESCE(ic.ref_count, 0) as ref_count
                FROM {self.table_name} i
                LEFT JOIN image_ppt_mapping m ON i.img_hash = m.img_hash
                LEFT JOIN image_counts ic ON i.img_hash = ic.img_hash
                ORDER BY COALESCE(ic.ref_count, 0) DESC, i.extract_date DESC
            """
            
            # 只在指定了limit时添加分页
            if limit is not None:
                query += " LIMIT ? OFFSET ?"
                self.cursor.execute(query, [limit, offset])
            else:
                self.cursor.execute(query)
            
            results = self.cursor.fetchall()
            
            images = []
            for row in results:
                (img_hash, img_path, img_name, extract_date, img_type, 
                 width, height, ppt_path, ref_count) = row
                if os.path.exists(img_path):  # 只返回仍然存在的图片
                    images.append({
                        'hash': img_hash,
                        'path': img_path,
                        'name': img_name,
                        'extract_date': extract_date,
                        'type': img_type,
                        'width': width,
                        'height': height,
                        'ppt_path': ppt_path,
                        'ppt_name': Path(ppt_path).name if ppt_path else None,
                        'ref_count': ref_count
                    })
            
            return images
            
        except Exception as e:
            print(f"获取图片列表失败: {str(e)}")
            return []

    def search_images(self, keyword: str = "", img_type: str = None) -> list:
        """搜索图片"""
        try:
            query = f"""
                SELECT DISTINCT i.img_hash, i.img_path, i.img_name, i.extract_date, 
                       i.img_type, i.width, i.height,
                       m.pptx_path
                FROM {self.table_name} i
                LEFT JOIN image_ppt_mapping m ON i.img_hash = m.img_hash
                WHERE 1=1
            """
            params = []
            
            if keyword:
                query += """ AND (
                    i.img_name LIKE ? OR 
                    i.img_path LIKE ? OR 
                    m.pptx_path LIKE ?
                )"""
                params.extend([f"%{keyword}%"] * 3)
            
            if img_type:
                query += " AND i.img_type = ?"
                params.append(img_type)
            
            query += " ORDER BY i.extract_date DESC"
            
            self.cursor.execute(query, params)
            results = self.cursor.fetchall()
            
            images = []
            for row in results:
                img_hash, img_path, img_name, extract_date, img_type, width, height, ppt_path = row
                if os.path.exists(img_path):
                    images.append({
                        'hash': img_hash,
                        'path': img_path,
                        'name': img_name,
                        'extract_date': extract_date,
                        'type': img_type,
                        'width': width,
                        'height': height,
                        'ppt_path': ppt_path,
                        'ppt_name': Path(ppt_path).name if ppt_path else None
                    })
            return images
        except Exception as e:
            print(f"搜索图片时出错: {str(e)}")
            return []

    def get_image_info(self, img_hash: str) -> dict:
        """获取指定图片的详信息"""
        if not self.cursor:
            raise ValueError("数据库未初始化")
        
        try:
            self.cursor.execute(
                f"SELECT img_path, pptx_path, extract_date FROM {self.table_name} "
                f"WHERE img_hash = ?",
                (img_hash,)
            )
            result = self.cursor.fetchone()
            
            if result:
                img_path, pptx_path, extract_date = result
                return {
                    'hash': img_hash,
                    'path': img_path,
                    'ppt_path': pptx_path,
                    'extract_date': extract_date,
                    'name': Path(img_path).name,
                    'ppt_name': Path(pptx_path).name
                }
            return None
        except Exception as e:
            print(f"获取图片信息时出错: {str(e)}")
            return None

    def get_image_stats(self) -> dict:
        """获取图片库统计信息"""
        try:
            stats = {
                'total': 0,
                'background': 0,
                'icon': 0,
                'normal': 0,
                'ppt_count': 0
            }
            
            # 获取图片统计
            self.cursor.execute(f"""
                SELECT img_type, COUNT(*) 
                FROM {self.table_name}
                GROUP BY img_type
            """)
            for img_type, count in self.cursor.fetchall():
                stats[img_type] = count
                stats['total'] += count
            
            # 获取PPT数量
            self.cursor.execute("""
                SELECT COUNT(DISTINCT pptx_path) 
                FROM image_ppt_mapping
            """)
            stats['ppt_count'] = self.cursor.fetchone()[0]
            
            return stats
        except Exception as e:
            print(f"获取统计信息时出错: {str(e)}")
            return {}

    def save_setting(self, key: str, value: str):
        """保存设置"""
        self.cursor.execute(
            "INSERT OR REPLACE INTO settings (key, value) VALUES (?, ?)",
            (key, value)
        )
        self.db_conn.commit()

    def get_setting(self, key: str, default: str = None) -> str:
        """获取设置"""
        self.cursor.execute("SELECT value FROM settings WHERE key = ?", (key,))
        result = self.cursor.fetchone()
        return result[0] if result else default

    def add_ppt_source(self, path: str):
        """添加PPT源文件夹"""
        self.cursor.execute(
            "INSERT OR REPLACE INTO ppt_sources (path, added_date) VALUES (?, ?)",
            (path, datetime.now().strftime('%Y-%m-%d %H:%M:%S'))
        )
        self.db_conn.commit()

    def get_ppt_sources(self) -> list:
        """获取所有PPT源文件夹"""
        self.cursor.execute("SELECT path FROM ppt_sources")
        return [row[0] for row in self.cursor.fetchall()]

    def get_image_ppt_mappings(self, img_hash: str) -> list:
        """获取图片被使用的所有PPT信息"""
        self.cursor.execute(
            """
            SELECT m.pptx_path, m.slide_index, m.shape_index, m.extract_date
            FROM image_ppt_mapping m
            WHERE m.img_hash = ?
            ORDER BY m.extract_date DESC
            """,
            (img_hash,)
        )
        
        mappings = []
        for row in self.cursor.fetchall():
            ppt_path, slide_idx, shape_idx, extract_date = row
            if os.path.exists(ppt_path):  # 只返回存在的PPT
                mappings.append({
                    'ppt_path': ppt_path,
                    'ppt_name': Path(ppt_path).name,
                    'slide': slide_idx,
                    'shape': shape_idx,
                    'extract_date': extract_date
                })
        return mappings

    @lru_cache(maxsize=1000)
    def _get_cached_thumbnail(self, img_path: str, ref_count: int) -> str:
        """获取缓存的缩略图路径"""
        return self._create_thumbnail_with_badge(img_path, ref_count)

    def _create_thumbnail_with_badge(self, img_path: str, ref_count: int, size=(200, 200)) -> str:
        """创建带角标的缩略图并缓存"""
        try:
            # 生成缓存文件名
            cache_dir = Path(self._get_app_data_dir()) / "thumbnails"
            cache_dir.mkdir(parents=True, exist_ok=True)
            
            # 使用图片路径和引用次数生成缓存文件名
            thumb_hash = hashlib.md5(f"{img_path}_{ref_count}".encode()).hexdigest()
            cache_path = cache_dir / f"{thumb_hash}.png"
            
            # 如果缓存存在且新于原图，直接返回
            if cache_path.exists() and cache_path.stat().st_mtime > Path(img_path).stat().st_mtime:
                return str(cache_path)
            
            # 创建缩略图
            with Image.open(img_path) as img:
                # 保持RGBA格式
                if img.mode == 'RGBA':
                    thumb = img.copy()
                    thumb.thumbnail(size, Image.Resampling.LANCZOS)
                else:
                    thumb = img.convert('RGBA')
                    thumb.thumbnail(size, Image.Resampling.LANCZOS)
                
                # 如果有引用计数，添加角标
                if ref_count > 1:
                    # 创建一个新的透明图层用于绘角标
                    badge = Image.new('RGBA', thumb.size, (0, 0, 0, 0))
                    draw = ImageDraw.Draw(badge)
                    
                    # 绘制圆形背景
                    badge_size = 24
                    x = thumb.size[0] - badge_size - 5
                    y = 5
                    draw.ellipse(
                        [x, y, x + badge_size, y + badge_size],
                        fill=(255, 87, 34, 255)  # 橙色
                    )
                    
                    # 添加文字
                    try:
                        font = ImageFont.truetype("arial.ttf", 12)
                    except:
                        # 如果找不到 arial.ttf，使用默认字体
                        font = ImageFont.load_default()
                        
                    text = str(ref_count)
                    text_bbox = draw.textbbox((0, 0), text, font=font)
                    text_width = text_bbox[2] - text_bbox[0]
                    text_height = text_bbox[3] - text_bbox[1]
                    text_x = x + (badge_size - text_width) // 2
                    text_y = y + (badge_size - text_height) // 2
                    draw.text((text_x, text_y), text, fill=(255, 255, 255, 255), font=font)
                    
                    # 合并图层
                    thumb = Image.alpha_composite(thumb, badge)
                
                # 保存缩略图，保持透明通道
                thumb.save(cache_path, "PNG", optimize=True)
                return str(cache_path)
            
        except Exception as e:
            print(f"创建缩略图失败: {str(e)}")
            return img_path
    
    def _init_clip_model(self):
        """延迟载CLIP模型"""
        if not self.clip_available:
            raise RuntimeError("CLIP模型相关库未安装，无法使用标签识别功能")
        
        if self.clip_model is None:
            try:
                # 设置镜像站点
                os.environ['HF_ENDPOINT'] = 'https://hf-mirror.com'
                os.environ['HF_HOME'] = str(Path(self._get_app_data_dir()) / "huggingface")
                os.environ['HF_MIRROR'] = 'https://hf-mirror.com'
                
                # 使用本地模型
                model_name = "openai/clip-vit-base-patch32"
                
                # 创建缓存目录
                cache_dir = Path(self._get_app_data_dir()) / "models" / "clip"
                cache_dir.mkdir(parents=True, exist_ok=True)
                
                # 忽略警告
                import warnings
                warnings.filterwarnings("ignore", message=".*resume_download.*")
                
                # 加载模型和处理器
                print("正在加载CLIP模型和处理器...")
                from transformers import CLIPModel, CLIPProcessor
                
                # 尝试从本地加载
                try:
                    print(f"尝试从本地加载模型: {cache_dir}")
                    self.clip_model = CLIPModel.from_pretrained(
                        model_name,
                        cache_dir=str(cache_dir),
                        local_files_only=True,
                        mirror='tuna'  # 使用清华镜像
                    )
                    self.clip_processor = CLIPProcessor.from_pretrained(
                        model_name,
                        cache_dir=str(cache_dir),
                        local_files_only=True,
                        mirror='tuna'
                    )
                    print("从本地加载模型成功")
                except Exception as local_e:
                    print(f"本地模型不存在: {str(local_e)}")
                    print("尝试从网络下载...")
                    
                    # 如果地载失败，尝试从网络下载
                    print("使用镜像下载模型...")
                    self.clip_model = CLIPModel.from_pretrained(
                        model_name,
                        cache_dir=str(cache_dir),
                        local_files_only=False,
                        resume_download=True,
                        mirror='tuna',
                        proxies={'http': 'http://mirrors.aliyun.com/pypi/simple/'}
                    )
                    
                    print("使用镜像下载处理器...")
                    self.clip_processor = CLIPProcessor.from_pretrained(
                        model_name,
                        cache_dir=str(cache_dir),
                        local_files_only=False,
                        resume_download=True,
                        mirror='tuna',
                        proxies={'http': 'http://mirrors.aliyun.com/pypi/simple/'}
                    )
                    print("从网络下载模型成功")
                
                # 移动到GPU（如果可用）
                if torch.cuda.is_available():
                    self.clip_model = self.clip_model.to('cuda')
                    print("模型已加载到GPU")
                else:
                    print("模型将在CPU上运行")
                
                # 设置为评估模式
                self.clip_model.eval()
                print("CLIP模型加载完成")
                
            except Exception as e:
                print(f"加载CLIP模型失败: {str(e)}")
                print("\n尝试使用以下方法解决:")
                print("1. 检查网络连接")
                print("2. 使用代理或VPN")
                print("3. 手动载模型文件:", cache_dir)
                print("4. 或者手动下载模型文件：")
                print("   - 访问：https://huggingface.co/openai/clip-vit-base-patch32")
                print("   - 下载所需文件并放到:", cache_dir)
                print("\n需要下载的文件:")
                print("- config.json")
                print("- pytorch_model.bin")
                print("- vocab.json")
                print("- merges.txt")
                self.clip_model = None
                self.clip_processor = None
                raise
    
    def add_tags(self, tags: List[Dict[str, str]]):
        """添加标签到数库"""
        try:
            for tag in tags:
                # 处理标签文本
                name = tag['name'].strip().lower()  # 标准化标签名称
                
                # 移除特殊字符，只保留英文、数字和空格
                import re
                name = re.sub(r'[^a-zA-Z0-9\s-]', '', name)
                
                self.cursor.execute(
                    "INSERT OR IGNORE INTO tags (name, category, created_at) VALUES (?, ?, ?)",
                    (name, tag.get('category', 'uncategorized'), 
                     datetime.now().strftime('%Y-%m-%d %H:%M:%S'))
                )
            self.db_conn.commit()
        except Exception as e:
            print(f"添加标签失败: {str(e)}")
            self.db_conn.rollback()
    
    def process_image_tags(self, img_hash: str, confidence_threshold: float = 0.5):
        """处理单个图片的标签"""
        try:
            # 检查CLIP模型和处理器
            if not self.clip_available:
                print("CLIP功能未启用")
                return
            
            if self.clip_model is None or self.clip_processor is None:
                print("CLIP模型未初始化，尝试初始化...")
                self._init_clip_model()
                
                if self.clip_model is None or self.clip_processor is None:
                    print("CLIP模型初始化失败")
                    return
            
            # 获取图片路径
            self.cursor.execute(
                f"SELECT img_path, img_name FROM {self.table_name} WHERE img_hash = ?",
                (img_hash,)
            )
            result = self.cursor.fetchone()
            if not result:
                print(f"找不到图片记录: {img_hash}")
                return
            
            img_path, img_name = result
            if not os.path.exists(img_path):
                print(f"图片文件不存在: {img_path}")
                return
            
            print(f"\n处理图片: {img_name}")
            
            # 获取所有标签
            self.cursor.execute("SELECT id, name FROM tags")
            tags = self.cursor.fetchall()
            if not tags:
                print("没有可用的标签")
                return
            
            # 准备图片和文本
            print("加载图片...")
            try:
                image = Image.open(img_path)
                
                # 图像预处理
                print("预处理图像...")
                # 1. 确保图像是RGB模式
                if image.mode not in ['RGB']:
                    print(f"转换图片格式从 {image.mode} 到 RGB")
                    image = image.convert('RGB')
                
                # 2. 调整图像大小到合适的尺寸
                target_size = (224, 224)  # CLIP模型的标准输入尺寸
                if image.size != target_size:
                    print(f"调整图像尺寸从 {image.size} 到 {target_size}")
                    image = image.resize(target_size, Image.Resampling.LANCZOS)
                
                # 3. 标准化图像亮度和对比度
                from PIL import ImageEnhance
                enhancer = ImageEnhance.Brightness(image)
                image = enhancer.enhance(1.1)  # 略微提高亮度
                
                enhancer = ImageEnhance.Contrast(image)
                image = enhancer.enhance(1.1)  # 略微提高对比度
                
                # 4. 添加图像增强
                enhancer = ImageEnhance.Sharpness(image)
                image = enhancer.enhance(1.2)  # 适度锐化
                
                print("图像预处理完成")
                
            except Exception as e:
                print(f"加载或处理图片失败: {str(e)}")
                return
            
            # 准备标签文本
            tag_texts = [tag[1] for tag in tags]
            print(f"处理 {len(tag_texts)} 个标签...")
            
            # 使用CLIP进行预测
            print("进行标签识别...")
            try:
                # 对每个标签添加提示词以提高准确性
                enhanced_texts = []
                for text in tag_texts:
                    prompts = [
                        f"a photo of {text}",
                        f"an image of {text}",
                        f"a picture showing {text}",
                        text
                    ]
                    enhanced_texts.extend(prompts)
                
                # 准备输入
                inputs = self.clip_processor(
                    images=image,
                    text=enhanced_texts,
                    return_tensors="pt",
                    padding=True,
                    max_length=77,
                    truncation=True
                )
                
                # 移动到GPU
                if torch.cuda.is_available():
                    inputs = {k: v.to('cuda') for k, v in inputs.items()}
                    print("已将数据移动到GPU")
                else:
                    print("在CPU上运行推理")
                
                # 获取预测结果
                with torch.no_grad():
                    outputs = self.clip_model(**inputs)
                    logits_per_image = outputs.logits_per_image
                    
                    # 1. 计算 cosine similarity
                    logits_per_image = logits_per_image / logits_per_image.norm(dim=-1, keepdim=True)
                    
                    # 2. 使用更低的温度系数来增加区分度
                    temperature = 0.07
                    logits_per_image = logits_per_image / temperature
                    
                    # 3. 计算每个提示词的概率
                    probs = torch.softmax(logits_per_image, dim=-1)  # 使用 softmax 而不是 sigmoid
                    probs = probs.cpu().numpy()[0]
                    
                    # 4. 将多个提示词的结果合并回原始标签
                    num_prompts = 4
                    tag_probs = []
                    for i in range(0, len(probs), num_prompts):
                        prompt_probs = probs[i:i + num_prompts]
                        # 使用最大值而不是平均值
                        tag_prob = max(prompt_probs)
                        tag_probs.append(tag_prob)
                    
                    # 5. 再次进行 softmax 归一化，增加区分度
                    tag_probs = torch.softmax(torch.tensor(tag_probs) / temperature, dim=0).numpy()
                    
                    # 6. 调整阈值
                    confidence_threshold = 0.1  # 降低阈值，因为概率和为1
                    
                    print("\n标签概率分布:")
                    for (_, tag_name), prob in zip(tags, tag_probs):
                        print(f"{tag_name}: {prob:.3f}")
                
                # 先删除旧的标签
                print("\n删除旧标签...")
                self.cursor.execute(
                    "DELETE FROM image_tags WHERE img_hash = ?",
                    (img_hash,)
                )
                
                # 保存高于阈值的标签
                matched_tags = []
                for (tag_id, tag_name), prob in zip(tags, tag_probs):
                    if prob >= confidence_threshold:  # 只要概率大于阈值就添加标签
                        self.cursor.execute(
                            "INSERT INTO image_tags "
                            "(img_hash, tag_id, confidence, created_at) VALUES (?, ?, ?, ?)",
                            (img_hash, tag_id, float(prob), 
                             datetime.now().strftime('%Y-%m-%d %H:%M:%S'))
                        )
                        matched_tags.append(f"{tag_name}({prob:.3f})")
                
                if matched_tags:
                    print(f"\n匹配到的标签: {', '.join(matched_tags)}")
                else:
                    print("\n未匹配到任何标签")
                
                self.db_conn.commit()
                print("标签数据已更新")
                
            except Exception as e:
                print(f"CLIP推理失败: {str(e)}")
                self.db_conn.rollback()
            
        except Exception as e:
            print(f"处理图片标签失败: {str(e)}")
            if hasattr(e, '__traceback__'):
                import traceback
                traceback.print_tb(e.__traceback__)
    
    def get_image_tags(self, img_hash: str) -> List[Dict]:
        """获取图片的标签"""
        try:
            self.cursor.execute(
                """
                SELECT t.name, t.category, it.confidence
                FROM image_tags it
                JOIN tags t ON it.tag_id = t.id
                WHERE it.img_hash = ?
                ORDER BY it.confidence DESC
                """,
                (img_hash,)
            )
            
            return [{
                'name': name,
                'category': category,
                'confidence': confidence
            } for name, category, confidence in self.cursor.fetchall()]
            
        except Exception as e:
            print(f"获取图片标签失败: {str(e)}")
            return []

    def get_all_tags(self) -> List[Dict]:
        """获取所有标签"""
        try:
            self.cursor.execute(
                """
                SELECT t.id, t.name, t.category, t.created_at,
                       COUNT(it.img_hash) as usage_count
                FROM tags t
                LEFT JOIN image_tags it ON t.id = it.tag_id
                GROUP BY t.id
                ORDER BY t.category, t.name
                """
            )
            
            return [{
                'id': tag_id,
                'name': name,
                'category': category,
                'created_at': created_at,
                'usage_count': usage_count
            } for tag_id, name, category, created_at, usage_count in self.cursor.fetchall()]
            
        except Exception as e:
            print(f"获取标签列表失败: {str(e)}")
            return []

    def get_tag_categories(self) -> List[str]:
        """获取所有标签类别"""
        try:
            self.cursor.execute("SELECT DISTINCT category FROM tags ORDER BY category")
            return [row[0] for row in self.cursor.fetchall()]
        except Exception as e:
            print(f"获取标签类别失败: {str(e)}")
            return []

    @lru_cache(maxsize=128)
    def search_images_by_tags(self, tag_names: tuple, match_all: bool = True) -> List[Dict]:
        """根据标签搜索图片（使用元组作为参数以支持缓存）"""
        try:
            if not tag_names:
                return []
            
            # 将标签名转换为小写
            tag_names = [name.lower() for name in tag_names]
            
            # 构建优化后的查询
            query = f"""
                WITH matched_tags AS (
                    -- 找到所有匹配的标签ID
                    SELECT id, name
                    FROM tags
                    WHERE LOWER(name) IN ({','.join(['?' for _ in tag_names])})
                ),
                tagged_images AS (
                    -- 找到包含这些标签的图片
                    SELECT 
                        it.img_hash,
                        COUNT(DISTINCT mt.name) as matched_tag_count,
                        GROUP_CONCAT(mt.name || '(' || ROUND(it.confidence, 2) || ')') as matched_tags
                    FROM image_tags it
                    JOIN matched_tags mt ON it.tag_id = mt.id
                    GROUP BY it.img_hash
                    {f'HAVING matched_tag_count = {len(tag_names)}' if match_all else ''}
                ),
                image_counts AS (
                    SELECT img_hash, COUNT(DISTINCT pptx_path) as ref_count
                    FROM image_ppt_mapping
                    GROUP BY img_hash
                )
                SELECT 
                    i.img_hash, i.img_path, i.img_name, i.extract_date, 
                    i.img_type, i.width, i.height,
                    m.pptx_path,
                    COALESCE(ic.ref_count, 0) as ref_count,
                    ti.matched_tags,
                    ti.matched_tag_count
                FROM tagged_images ti
                JOIN {self.table_name} i ON ti.img_hash = i.img_hash
                LEFT JOIN image_ppt_mapping m ON i.img_hash = m.img_hash
                LEFT JOIN image_counts ic ON i.img_hash = ic.img_hash
                ORDER BY 
                    ti.matched_tag_count DESC,
                    ic.ref_count DESC
            """
            
            print(f"执行标签搜索: {tag_names}")
            print(f"SQL查询: {query}")
            
            self.cursor.execute(query, tag_names)
            results = self.cursor.fetchall()
            
            images = []
            for row in results:
                (img_hash, img_path, img_name, extract_date, img_type, 
                 width, height, ppt_path, ref_count, tags, matched_count) = row
                
                if os.path.exists(img_path):
                    images.append({
                        'hash': img_hash,
                        'path': img_path,
                        'name': img_name,
                        'extract_date': extract_date,
                        'type': img_type,
                        'width': width,
                        'height': height,
                        'ppt_path': ppt_path,
                        'ppt_name': Path(ppt_path).name if ppt_path else None,
                        'ref_count': ref_count,
                        'tags': tags.split(',') if tags else [],
                        'matched_tag_count': matched_count
                    })
            
            print(f"找到 {len(images)} 张匹配的图片")
            return images
            
        except Exception as e:
            print(f"按标签搜索图片失败: {str(e)}")
            print("查询参数:", tag_names)
            print("错误详情:", e.__class__.__name__)
            if hasattr(e, '__traceback__'):
                import traceback
                traceback.print_tb(e.__traceback__)
            return []

    def batch_process_tags(self, confidence_threshold: float = 0.5, 
                          progress_callback=None) -> int:
        """
        批量处理所有图片的标签
        回处理的图片数量
        """
        try:
            # 获取所有图片
            self.cursor.execute(f"SELECT img_hash FROM {self.table_name}")
            image_hashes = [row[0] for row in self.cursor.fetchall()]
            
            processed_count = 0
            total = len(image_hashes)
            
            for i, img_hash in enumerate(image_hashes):
                try:
                    # 更新进度
                    if progress_callback and callable(progress_callback):
                        should_continue = progress_callback(i + 1, total)
                        if should_continue is False:  # 允许取消处理
                            break
                    
                    # 处理标签
                    self.process_image_tags(img_hash, confidence_threshold)
                    processed_count += 1
                    
                except Exception as e:
                    print(f"处理图片 {img_hash} 标签失败: {str(e)}")
                    continue
            
            return processed_count
            
        except Exception as e:
            print(f"批量处理标签失败: {str(e)}")
            return 0

    def update_tag(self, tag_id: int, name: str, category: str = None):
        """更新标签"""
        try:
            self.cursor.execute(
                "UPDATE tags SET name = ?, category = ?, created_at = ? WHERE id = ?",
                (name, category, datetime.now().strftime('%Y-%m-%d %H:%M:%S'), tag_id)
            )
            self.db_conn.commit()
        except Exception as e:
            print(f"更新标签失败: {str(e)}")
            self.db_conn.rollback()

    def delete_tag(self, tag_id: int):
        """删除标签"""
        try:
            self.cursor.execute("DELETE FROM tags WHERE id = ?", (tag_id,))
            self.db_conn.commit()
        except Exception as e:
            print(f"删除标签失败: {str(e)}")
            self.db_conn.rollback()

    def batch_create_thumbnails(self, image_paths, ref_counts):
        """批量创建缩略图"""
        thumb_paths = []
        for img_path, ref_count in zip(image_paths, ref_counts):
            try:
                thumb_path = self._create_thumbnail_with_badge(img_path, ref_count)
                thumb_paths.append(thumb_path)
            except Exception as e:
                print(f"创建缩略图失败: {str(e)}")
                thumb_paths.append(None)
        return thumb_paths

    def add_tag_category(self, name: str):
        """添加标签分类"""
        try:
            self.cursor.execute(
                "INSERT INTO tag_categories (name, created_at) VALUES (?, ?)",
                (name, datetime.now().strftime('%Y-%m-%d %H:%M:%S'))
            )
            self.db_conn.commit()
        except Exception as e:
            print(f"添加标签分类失败: {str(e)}")
            self.db_conn.rollback()

    def update_tag_category(self, old_name: str, new_name: str):
        """更新标签分类"""
        try:
            self.cursor.execute(
                "UPDATE tags SET category = ? WHERE category = ?",
                (new_name, old_name)
            )
            self.db_conn.commit()
        except Exception as e:
            print(f"更新标签分类失败: {str(e)}")
            self.db_conn.rollback()

    def delete_tag_category(self, name: str):
        """删除标签分类"""
        try:
            # 将该分类下的标签移至"未分类"
            self.cursor.execute(
                "UPDATE tags SET category = NULL WHERE category = ?",
                (name,)
            )
            self.db_conn.commit()
        except Exception as e:
            print(f"删除标签分类失败: {str(e)}")
            self.db_conn.rollback()

    def get_tag_id(self, tag_name: str) -> int:
        """根据标签名获取标签ID"""
        try:
            self.cursor.execute(
                "SELECT id FROM tags WHERE name = ?",
                (tag_name,)
            )
            result = self.cursor.fetchone()
            return result[0] if result else None
        except Exception as e:
            print(f"获取标签ID失败: {str(e)}")
            return None